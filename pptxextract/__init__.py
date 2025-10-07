import io, os, re, json, logging
from typing import Any, Dict, List, Optional, Tuple

import requests
import azure.functions as func
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_PX = 9525  # ~96dpi

EMAIL_RE = re.compile(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}", re.I)
PHONE_RE = re.compile(r"(?:\+?\d[\d ()\-]{7,}\d)")
URL_RE   = re.compile(r"(?:(?:https?://)|(?:www\.))\S+", re.I)
LINKEDIN_RE = re.compile(r"(?:linkedin\.com/[A-Za-z0-9_\-/]+)", re.I)

def _px(emu): return int(emu / EMU_PER_PX)

def _download_pptx(sas: str) -> bytes:
    r = requests.get(sas, timeout=180)
    r.raise_for_status()
    return r.content

def _text_from_textframe(tf) -> Tuple[List[str], Optional[float]]:
    lines: List[str] = []
    max_pt = None
    for p in tf.paragraphs:
        t = "".join((run.text or "") + (f" ({run.hyperlink.address})" if getattr(run, "hyperlink", None) and run.hyperlink.address else "") for run in p.runs).strip()
        if not t: continue
        indent = "  " * (p.level or 0)
        lines.append(indent + t)
        for r in p.runs:
            try:
                if r.font and r.font.size:
                    v = float(r.font.size.pt)
                    max_pt = v if max_pt is None or v > max_pt else max_pt
            except Exception:
                pass
    return lines, max_pt

def _lines_from_table(tbl) -> List[str]:
    out: List[str] = []
    for r in tbl.rows:
        cells = []
        for c in r.cells:
            t = (c.text or "").strip()
            if t: t = re.sub(r"\s*\n\s*", " / ", t)
            cells.append(t)
        row_txt = " | ".join(filter(None, cells)).strip()
        if row_txt: out.append(row_txt)
    return out

def _visit_shape(sh, left_cut: int, acc: List[Dict[str, Any]]):
    lt, tp, w, h = _px(sh.left), _px(sh.top), _px(sh.width), _px(sh.height)
    col = "left" if lt + w/2 <= left_cut else "right"

    if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
        for g in sh.shapes:
            _visit_shape(g, left_cut, acc)
        return

    if getattr(sh, "has_text_frame", False) and sh.has_text_frame:
        lines, max_pt = _text_from_textframe(sh.text_frame)
        if lines:
            is_title = False
            try:
                if getattr(sh, "is_placeholder", False) and getattr(sh.placeholder_format, "type", None) == 1:
                    is_title = True
            except Exception:
                pass
            acc.append({"type":"text","col":col,"lines":lines,"bbox":[lt,tp,w,h],"is_title":is_title,"max_pt":max_pt})

    if getattr(sh, "has_table", False) and sh.has_table:
        lines = _lines_from_table(sh.table)
        if lines:
            acc.append({"type":"table","col":col,"lines":lines,"bbox":[lt,tp,w,h]})

    # alt text fallback
    try:
        alt = (sh.alternative_text or "").strip()
        if alt:
            acc.append({"type":"alt","col":col,"lines":[alt],"bbox":[lt,tp,w,h]})
    except Exception:
        pass

def _extract_slide(slide, slide_width_px: int) -> Dict[str, Any]:
    left_cut = int(slide_width_px * 0.45)
    blocks: List[Dict[str, Any]] = []
    for sh in slide.shapes:
        _visit_shape(sh, left_cut, blocks)

    blocks.sort(key=lambda b: (b["bbox"][1], b["bbox"][0]))

    # notes
    notes = None
    try:
        if getattr(slide, "notes_slide", None) and slide.notes_slide and slide.notes_slide.notes_text_frame:
            t = (slide.notes_slide.notes_text_frame.text or "").strip()
            if t: notes = t
    except Exception:
        pass

    # linear view (with L/R tags)
    linear: List[str] = []
    title_guess, title_pt = None, 0.0
    for b in blocks:
        pref = "[L] " if b.get("col")=="left" else "[R] "
        for ln in b.get("lines", []):
            linear.append(pref + ln)
        if b.get("type")=="text" and b.get("max_pt"):
            if b["max_pt"] > title_pt:
                title_pt = b["max_pt"]; title_guess = " ".join(b["lines"][:1])[:140]
    if notes: linear.append(f"Notes: {notes}")

    # explicit title
    title = None
    for b in blocks:
        if b.get("is_title"):
            title = " ".join([ln.strip() for ln in b["lines"][:1]])
            break
    if not title and title_guess: title = title_guess.strip()

    return {"title": title, "blocks": blocks, "notes": notes, "text": "\n".join(linear).strip()}

def _gather_hints(text: str) -> Dict[str, List[str]]:
    emails = EMAIL_RE.findall(text) or []
    phones = [p.strip() for p in PHONE_RE.findall(text) or []]
    urls   = URL_RE.findall(text) or []
    linked = LINKEDIN_RE.findall(text) or []
    return {
        "emails": sorted(set(emails)),
        "phones": sorted(set(phones)),
        "urls":   sorted(set(urls)),
        "linkedin": sorted(set(linked))
    }

def main(req: func.HttpRequest) -> func.HttpResponse:
    if req.method != "POST": return func.HttpResponse("POST only", status_code=405)
    try: body = req.get_json()
    except ValueError: return func.HttpResponse("Invalid JSON", status_code=400)

    sas = body.get("ppt_blob_sas")
    if not sas: return func.HttpResponse("Missing 'ppt_blob_sas'", status_code=400)

    try:
        data = _download_pptx(sas)
        prs = Presentation(io.BytesIO(data))
    except Exception as e:
        logging.exception("Failed to open PPTX")
        return func.HttpResponse(f"Could not read PPTX: {e}", status_code=400)

    slide_w_px = _px(prs.slide_width)
    slides: List[Dict[str, Any]] = [ _extract_slide(s, slide_w_px) for s in prs.slides ]

    parts: List[str] = []
    for i, sl in enumerate(slides, 1):
        if sl.get("title"): parts.append(f"[Slide {i}] {sl['title']}")
        if sl.get("text"):  parts.append(sl["text"])
    all_text = "\n\n".join(parts).strip()
    hints = _gather_hints(all_text)

    return func.HttpResponse(
        json.dumps({"ok": True, "slides": slides, "slides_text": all_text, "raw": all_text, "hints": hints}),
        status_code=200, mimetype="application/json")
